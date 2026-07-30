from __future__ import annotations

import argparse
from pathlib import Path
from typing import Dict, Iterable, List, Optional

import gspread
from google.oauth2.service_account import Credentials
from pptx import Presentation

import streamlit as st


SHEET_ID = "1tE4uXTSaO6rHk_W3GXU55FEwUfMw_CLdnMTXQJCWYx8"
WORKSHEET_NAME = "BPAD"

SCOPES = (
    "https://www.googleapis.com/auth/spreadsheets.readonly",
    "https://www.googleapis.com/auth/drive.readonly",
)

BASE_DIR = Path(__file__).resolve().parent.parent
DEFAULT_CREDS = BASE_DIR / "secrets" / "google_credentials.json"
DEFAULT_TEMPLATE_DIR = BASE_DIR / "templates"
DEFAULT_OUTPUT_DIR = BASE_DIR / "outputs"


# ---------------------------------------------------------------------
# Google Sheet cell map
# ---------------------------------------------------------------------
# These pointers come from the BPAD worksheet mapping you provided.
# Keep this section as the single source of truth for the report values.

SLIDE_6 = {
    "title": "Till date (16-07-2026)",
    "brief_recruited": {
        "families": "T8",
        "brief_only_individuals": "U8",
        "only_brief_b_plus_d_individuals": "U11",
    },
    "deep_recruited": {
        "families": "T7",
        "individuals": "U7",
    },
}

SLIDE_7 = {
    "title": "Till date",
    "rows": {
        "Deep Phase 2": {"families": "V7", "individuals": "W7"},
        "Deep Phase 3": {"families": "X7", "individuals": "Y7"},
        "Deep Phase 4": {"families": "Z7", "individuals": "AA7"},
        "Due for Phase 2": {"families": "V8", "individuals": "W8"},
        "Due for Phase 3": {"families": "X8", "individuals": "Y8"},
    },
}

SLIDE_8 = {
    "table_1": {
        "title": "Till date",
        "rows": {
            "Due for Phase 2": {"families": "V8", "individuals": "W8"},
            "Not yet due for F/U Phase 2": {"families": "V9", "individuals": "W9"},
            "dead_after_deep": {"label": "After deep", "individuals": "T25"},
            "dead_after_brief": {"label": "After brief", "individuals": "T24"},
        },
    },
    "table_2": {
        "title": "Till date",
        "rows": {
            "Due for Phase 3": {"families": "X8", "individuals": "Y8"},
            "Not yet due for F/U Phase 3": {"families": "X9", "individuals": "Y9"},
            "Dead": {"label": "", "individuals": "Y10"},
        },
    },
}

SLIDE_9 = {
    "title": "Till date",
    "rows": {
        "Contacted and willing": "T16",
        "Unwilling/ withdrew consent": "T17",
        "Unable to contact": "T18",
        "Not answering": "T19",
        "Other reasons": "T20",
        "Yet to contact": "T21",
        "Total": "T23",
    },
}

SLIDE_10 = {
    "title": "Till date",
    "rows": {
        "Contacted and willing": "V16",
        "Unwilling/ withdrew consent": "V17",
        "Unable to contact": "V18",
        "Not answering": "V19",
        "Other reasons": "V20",
        "Yet to contact": "V21",
        "Total": "V23",
    },
}


def find_template(template_dir: Path) -> Path:
    candidates = sorted(template_dir.glob("*.pptx"))
    if not candidates:
        raise FileNotFoundError(
            f"No .pptx template found in: {template_dir}. "
            "Put your template PPTX in the templates/ folder."
        )
    if len(candidates) > 1:
        raise FileExistsError(
            f"Multiple PPTX files found in {template_dir}: "
            + ", ".join(p.name for p in candidates)
            + ". Keep only one template or pass --template explicitly."
        )
    return candidates[0]


def authorize(creds_path: Path):
    # First try Streamlit secrets (production)
    try:
        credentials = Credentials.from_service_account_info(
            dict(st.secrets["gcp_service_account"]),
            scopes=SCOPES,
        )
        return gspread.authorize(credentials)

    except Exception:
        # Fallback to local JSON credentials
        if not creds_path.exists():
            raise FileNotFoundError(
                f"Credentials file not found: {creds_path} "
                "and no [gcp_service_account] found in Streamlit secrets."
            )

        credentials = Credentials.from_service_account_file(
            str(creds_path),
            scopes=SCOPES,
        )
        return gspread.authorize(credentials)
        
def open_worksheet(gc) :
    sh = gc.open_by_key(SHEET_ID)
    return sh.worksheet(WORKSHEET_NAME)


def cell(ws, address: str) -> str:
    value = ws.acell(address).value
    return "" if value is None else str(value)


def safe_int(value: str) -> int:
    value = "" if value is None else str(value).strip()
    if value in {"", "-", "NA", "N/A"}:
        return 0
    value = value.replace(",", "")
    try:
        return int(float(value))
    except ValueError:
        return 0


def table_sum(*values: str) -> str:
    return str(sum(safe_int(v) for v in values))


def get_table(slide, index: int = 0):
    tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]
    if index >= len(tables):
        raise IndexError(f"Slide has only {len(tables)} table(s); requested {index}.")
    return tables[index]


def set_text(cell_obj, text: str):
    # Replace the cell text while leaving the slide layout intact.
    cell_obj.text = "" if text is None else str(text)


def fill_slide_6(ws, slide):
    table = get_table(slide, 0)

    # Title
    set_text(table.cell(0, 0), SLIDE_6["title"])

    # Brief recruited row:
    # Display families + two individual counts exactly as requested.
    brief_families = cell(ws, SLIDE_6["brief_recruited"]["families"])
    brief_only_individuals = cell(ws, SLIDE_6["brief_recruited"]["brief_only_individuals"])
    brief_b_plus_d = cell(ws, SLIDE_6["brief_recruited"]["only_brief_b_plus_d_individuals"])
    set_text(table.cell(2, 0), "Brief recruited")
    set_text(table.cell(2, 1), brief_families)
    set_text(
        table.cell(2, 2),
        f"{brief_only_individuals}(Only Brief)+{brief_b_plus_d}(Only Brief (B+D))",
    )

    # Deep recruited row
    set_text(table.cell(3, 0), "Deep recruited")
    set_text(table.cell(3, 1), cell(ws, SLIDE_6["deep_recruited"]["families"]))
    set_text(table.cell(3, 2), cell(ws, SLIDE_6["deep_recruited"]["individuals"]))


def fill_slide_7(ws, slide):
    table = get_table(slide, 0)
    set_text(table.cell(0, 0), SLIDE_7["title"])

    row_map = {
        2: "Deep Phase 2",
        3: "Deep Phase 3",
        4: "Deep Phase 4",
        5: "Due for Phase 2",
        6: "Due for Phase 3",
    }

    for row_idx, label in row_map.items():
        set_text(table.cell(row_idx, 0), label)
        set_text(table.cell(row_idx, 1), cell(ws, SLIDE_7["rows"][label]["families"]))
        set_text(table.cell(row_idx, 2), cell(ws, SLIDE_7["rows"][label]["individuals"]))


def fill_slide_8(ws, slide):
    tables = [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]
    if len(tables) < 2:
        raise RuntimeError("Slide 8 must contain two tables.")

    # Table 1
    table = tables[0]
    set_text(table.cell(0, 0), SLIDE_8["table_1"]["title"])

    set_text(table.cell(2, 0), "Due for Phase 2")
    set_text(table.cell(2, 1), cell(ws, SLIDE_8["table_1"]["rows"]["Due for Phase 2"]["families"]))
    set_text(table.cell(2, 2), cell(ws, SLIDE_8["table_1"]["rows"]["Due for Phase 2"]["individuals"]))

    set_text(table.cell(3, 0), "Not yet due for F/U Phase 2")
    set_text(table.cell(3, 1), cell(ws, SLIDE_8["table_1"]["rows"]["Not yet due for F/U Phase 2"]["families"]))
    set_text(table.cell(3, 2), cell(ws, SLIDE_8["table_1"]["rows"]["Not yet due for F/U Phase 2"]["individuals"]))

    set_text(table.cell(4, 0), "dead")
    set_text(table.cell(4, 1), "After deep")
    set_text(table.cell(4, 2), cell(ws, SLIDE_8["table_1"]["rows"]["dead_after_deep"]["individuals"]))

    set_text(table.cell(5, 0), "")
    set_text(table.cell(5, 1), "After brief")
    set_text(table.cell(5, 2), cell(ws, SLIDE_8["table_1"]["rows"]["dead_after_brief"]["individuals"]))

    total_1 = table_sum(
        cell(ws, SLIDE_8["table_1"]["rows"]["Due for Phase 2"]["individuals"]),
        cell(ws, SLIDE_8["table_1"]["rows"]["Not yet due for F/U Phase 2"]["individuals"]),
        cell(ws, SLIDE_8["table_1"]["rows"]["dead_after_deep"]["individuals"]),
        cell(ws, SLIDE_8["table_1"]["rows"]["dead_after_brief"]["individuals"]),
    )
    set_text(table.cell(6, 0), "Total")
    set_text(table.cell(6, 1), "")
    set_text(table.cell(6, 2), total_1)

    # Table 2
    table = tables[1]
    set_text(table.cell(0, 0), SLIDE_8["table_2"]["title"])

    set_text(table.cell(2, 0), "Due for Phase 3")
    set_text(table.cell(2, 1), cell(ws, SLIDE_8["table_2"]["rows"]["Due for Phase 3"]["families"]))
    set_text(table.cell(2, 2), cell(ws, SLIDE_8["table_2"]["rows"]["Due for Phase 3"]["individuals"]))

    set_text(table.cell(3, 0), "Not yet due for F/U Phase 3")
    set_text(table.cell(3, 1), cell(ws, SLIDE_8["table_2"]["rows"]["Not yet due for F/U Phase 3"]["families"]))
    set_text(table.cell(3, 2), cell(ws, SLIDE_8["table_2"]["rows"]["Not yet due for F/U Phase 3"]["individuals"]))

    set_text(table.cell(4, 0), "Dead")
    set_text(table.cell(4, 1), "")
    set_text(table.cell(4, 2), cell(ws, SLIDE_8["table_2"]["rows"]["Dead"]["individuals"]))

    total_2 = table_sum(
        cell(ws, SLIDE_8["table_2"]["rows"]["Due for Phase 3"]["individuals"]),
        cell(ws, SLIDE_8["table_2"]["rows"]["Not yet due for F/U Phase 3"]["individuals"]),
        cell(ws, SLIDE_8["table_2"]["rows"]["Dead"]["individuals"]),
    )
    set_text(table.cell(5, 0), "Total")
    set_text(table.cell(5, 1), "")
    set_text(table.cell(5, 2), total_2)


def fill_slide_9(ws, slide):
    table = get_table(slide, 0)
    set_text(table.cell(0, 0), SLIDE_9["title"])

    row_map = {
        2: "Contacted and willing",
        3: "Unwilling/ withdrew consent",
        4: "Unable to contact",
        5: "Not answering",
        6: "Other reasons",
        7: "Yet to contact",
        8: "Total",
    }

    for row_idx, label in row_map.items():
        set_text(table.cell(row_idx, 0), label)
        set_text(table.cell(row_idx, 1), cell(ws, SLIDE_9["rows"][label]))


def fill_slide_10(ws, slide):
    table = get_table(slide, 0)
    set_text(table.cell(0, 0), SLIDE_10["title"])

    row_map = {
        2: "Contacted and willing",
        3: "Unwilling/ withdrew consent",
        4: "Unable to contact",
        5: "Not answering",
        6: "Other reasons",
        7: "Yet to contact",
        8: "Total",
    }

    for row_idx, label in row_map.items():
        set_text(table.cell(row_idx, 0), label)
        set_text(table.cell(row_idx, 1), cell(ws, SLIDE_10["rows"][label]))


def generate_ppt(
    template_path: Path,
    output_path: Path,
    creds_path: Path = DEFAULT_CREDS,
) -> Path:
    gc = authorize(creds_path)
    ws = open_worksheet(gc)

    prs = Presentation(str(template_path))

    fill_slide_6(ws, prs.slides[5])
    fill_slide_7(ws, prs.slides[6])
    fill_slide_8(ws, prs.slides[7])
    fill_slide_9(ws, prs.slides[8])
    fill_slide_10(ws, prs.slides[9])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def parse_args():
    parser = argparse.ArgumentParser(
        description="Fill BPAD report slides 6-10 from Google Sheets data."
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=None,
        help="Path to the template PPTX. If omitted, the only .pptx in templates/ is used.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT_DIR / "bpad_report_filled.pptx",
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--creds",
        type=Path,
        default=DEFAULT_CREDS,
        help="Path to Google service account JSON.",
    )
    return parser.parse_args()


def main():
    args = parse_args()

    template = args.template
    if template is None:
        template = find_template(DEFAULT_TEMPLATE_DIR)

    if not template.exists():
        raise FileNotFoundError(f"Template not found: {template}")

    output = generate_ppt(
        template_path=template,
        output_path=args.output,
        creds_path=args.creds,
    )
    print(f"Saved: {output.resolve()}")


if __name__ == "__main__":
    main()
